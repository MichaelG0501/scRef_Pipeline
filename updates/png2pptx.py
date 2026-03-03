#!/usr/bin/env python3
"""Convert PDF slide PNGs to a .pptx file for a given update cycle.

Usage:
    python updates/png2pptx.py DDmon
    # e.g. python updates/png2pptx.py 02mar

Reads:  updates/DDmon/pptx_slides_png/slide-*.png
Writes: updates/DDmon/update_DDmon.pptx
"""
import sys
import glob
import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

if len(sys.argv) < 2:
    print("Usage: python png2pptx.py DDmon")
    print("  e.g. python png2pptx.py 02mar")
    sys.exit(1)

update_id = sys.argv[1]  # e.g., "02mar"

# Resolve paths relative to this script's location (updates/)
script_dir = os.path.dirname(os.path.abspath(__file__))
update_dir = os.path.join(script_dir, update_id)
png_dir = os.path.join(update_dir, "pptx_slides_png")
png_glob = os.path.join(png_dir, "slide-*.png")
out_pptx = os.path.join(update_dir, f"update_{update_id}.pptx")

# 16:9 widescreen
slide_w_in = 13.333
slide_h_in = 7.5

prs = Presentation()
prs.slide_width = Inches(slide_w_in)
prs.slide_height = Inches(slide_h_in)
blank_layout = prs.slide_layouts[6]

pngs = sorted(
    glob.glob(png_glob),
    key=lambda p: int(os.path.splitext(os.path.basename(p))[0].split("-")[-1]),
)

if not pngs:
    raise SystemExit(f"No PNGs found matching: {png_glob}")

for p in pngs:
    slide = prs.slides.add_slide(blank_layout)

    with Image.open(p) as im:
        w_px, h_px = im.size
        if w_px < 1500 or h_px < 800:
            print(f"Warning: {p} looks low-res: {w_px}x{h_px}px")

    slide.shapes.add_picture(
        p, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height
    )

prs.save(out_pptx)
print(f"Saved: {out_pptx} ({len(pngs)} slides)")
